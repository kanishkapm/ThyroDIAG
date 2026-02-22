import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
PROJECT_DIR = "/home/kanishka/.gemini/antigravity/playground/interstellar-plasma/thyroid_project2.0"
RESULTS_DIR = os.path.join(PROJECT_DIR, "results")
os.makedirs(RESULTS_DIR, exist_ok=True)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
TITLE_COLOR = RGBColor(44, 62, 80)
ACCENT_COLOR = RGBColor(52, 152, 219)
SUCCESS_COLOR = RGBColor(39, 174, 96)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR
    title_p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    sub_frame = sub_box.text_frame
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = subtitle
    sub_p.font.size = Pt(24)
    sub_p.font.color.rgb = ACCENT_COLOR
    sub_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(items):
        p = content_frame.paragraphs[0] if i == 0 else content_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.space_after = Pt(12)

def add_image_slide(prs, title, image_path, caption=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), width=Inches(12.333))
    else:
        p_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
        p_box.text_frame.text = f"[Image missing: {os.path.basename(image_path)}]"
        p_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
        cap_p = cap_box.text_frame.paragraphs[0]
        cap_p.text = caption
        cap_p.font.size = Pt(14)
        cap_p.font.italic = True
        cap_p.alignment = PP_ALIGN.CENTER

add_title_slide(prs, "ThyroDIAG: An Ensemble ML Approach for\nThyroid Cancer Subtype Classification", "Using Real Clinical Gene Expression Data\nKanishka P M | MSc Bioinformatics Project Review")

add_content_slide(prs, "1. Rationale / Background", [
    "Clinical Need: Accurate diagnosis of thyroid cancer subtypes (especially aggressive Anaplastic vs. indolent Papillary) is critical for treatment.",
    "The Problem with Current ML Approaches: Many rely on synthetic data, perform only binary classification, or suffer from genome-wide noise.",
    "Our Solution: A robust pipeline leveraging real clinical gene expression profiles and an advanced stacking ensemble for highly accurate, multi-class diagnostic support."
])

add_content_slide(prs, "2. Aim & Objective", [
    "Primary Aim: Develop & validate a high-accuracy, multi-class ensemble ML model classifying four thyroid tissue subtypes using NCBI GEO data.",
    "Objective 1: Integrate/preprocess raw series matrix data from GEO.",
    "Objective 2: Identify a clinically relevant biomarker panel using LASSO regression.",
    "Objective 3: Engineer a Stacking Ensemble combining diverse base models (RF, SVM, GB, KNN, MLP).",
    "Objective 4: Evaluate performance on the 4-class problem (Normal, PTC, FTC, ATC)."
])

add_content_slide(prs, "3. Materials and Methods", [
    "Data Source: Patient-derived microarray data from NCBI GEO (GSE33630, GSE60541, GSE53157).",
    "Dataset Size: 215 total samples across 4 classes.",
    "• Papillary (PTC): 114 (53.0%)",
    "• Normal: 70 (32.6%)",
    "• Anaplastic (ATC): 16 (7.4%)",
    "• Follicular (FTC): 15 (7.0%)",
    "Base Models: Random Forest, SVM (RBF), Gradient Boosting, KNN, MLP.",
    "Meta-Learner: Logistic Regression via 5-fold CV."
])

add_content_slide(prs, "4. Workflow", [
    "Step 1: Data Acquisition → Download & parse GEO series matrix files.",
    "Step 2: Preprocessing → Quality control, standard scaling, missing value imputation.",
    "Step 3: Feature Selection → LASSO to distill hundreds of genes into 61 core driver/biomarker genes.",
    "Step 4: Model Training → Train 5 base classifiers and the Stacking Ensemble meta-learner.",
    "Step 5: Validation → Evaluate using confusion matrices and cross-validation."
])

add_image_slide(prs, "Pipeline Visualization",
    os.path.join(RESULTS_DIR, "pipeline_workflow.png"),
    "End-to-end data processing and modeling workflow")

add_content_slide(prs, "5. Parameters and Metrics", [
    "Feature Selection: LASSO (Least Absolute Shrinkage and Selection Operator) to prevent overfitting.",
    "Evaluation Strategy: 5-fold stratified cross-validation for robustness across unbalanced classes.",
    "Performance Metrics Evaluated: Overall Accuracy, Precision, Recall (Sensitivity), and F1-Score.",
    "Evaluation Visualizations: Confusion Matrix, Feature Importance Plots."
])

add_content_slide(prs, "6. Results: Biomarker Discovery", [
    "Successfully identified 61 highly discriminative genes.",
    "Thyroid Differentiation: TG, TPO, DIO1, TSHR, NIS (SLC5A5)",
    "Cancer Driver Mutations: BRAF, RET, RAS family, PTEN, TP53",
    "EMT & Diagnostics: TERT, GALECTIN3, HBME1, KRT19, VIM",
    "Proliferation (High in ATC): MKI67, PCNA, TOP2A",
    "Immune Checkpoints: PDL1, CTLA4"
])

add_content_slide(prs, "7. Results: Model Performance", [
    "Individual Models: SVM achieved highest individual test accuracy (93.0%), MLP (90.7%).",
    "Final Ensemble Performance:",
    "• Test Accuracy: 93.0%",
    "• Cross-Validation Mean: 94.2%",
    "• Recall: 93.0%",
    "• F1-Score: 89.7%",
    "Multi-Class Success: Successfully differentiated the 4 subtypes despite class imbalances."
])

add_image_slide(prs, "Results: Confusion Matrix",
    os.path.join(RESULTS_DIR, "confusion_matrix.png"),
    "High accuracy classification across all four distinct tissue subtypes")

add_content_slide(prs, "8. Discussion & Conclusion", [
    "Clinical Significance: With 93% accuracy on real patient data, this model shows strong potential as a diagnostic support tool for pathologists.",
    "Biological Relevance: Extracted biomarker panel aligns directly with known thyroid biology.",
    "Future Directions: Wider clinical validation, direct histopathology workflow integration, and web app refinement."
])

add_content_slide(prs, "9. Novelty Points Summary", [
    "Real Clinical Data: Authentic patient-derived gene expression data, not synthetic.",
    "Robust Feature Selection: 61 biomarker genes via LASSO preventing the 'curse of dimensionality'.",
    "Advanced Stacking Ensemble: Combining 5 diverse models rather than a single classifier.",
    "Simultaneous Multi-Class: Accurately classifying 4 distinct thyroid tissue subtypes together.",
    "Translation-Ready: 93.02% accuracy with a complete end-to-end reproducible pipeline, including a live clinical web app."
])

add_content_slide(prs, "10. Key References", [
    "1. Edgar R, Domrachev M, Lash AE. Gene Expression Omnibus: NCBI gene expression and hybridization array data repository. Nucleic Acids Res. 2002;30(1):207-210.",
    "2. Cancer Genome Atlas Research Network. Integrated genomic characterization of papillary thyroid carcinoma. Cell. 2014;159(3):676-690.",
    "3. Xing M. BRAF mutation in thyroid cancer. Endocr Relat Cancer. 2005;12(2):245-262.",
    "4. Hastie T, Tibshirani R, Wainwright M. Statistical Learning with Sparsity: The Lasso and Generalizations. CRC Press; 2015."
])

add_title_slide(prs, "Thank You", "Live Web App Demo & Questions\nThyroDIAG v3.0")

output_path = os.path.join(RESULTS_DIR, "Thyroid_Cancer_Review_Presentation.pptx")
prs.save(output_path)
print(f"Presentation generated at: {output_path}")
